import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Set 16:9 widescreen size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    NAVY = RGBColor(43, 53, 96)
    ORANGE_RED = RGBColor(232, 70, 10)
    DEEP_NAVY = RGBColor(31, 45, 123)
    LAVENDER = RGBColor(232, 234, 246)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(68, 68, 68)

    def add_footer(slide):
        # 🏠 icon (using a shape or text)
        footer_y = prs.slide_height - Inches(0.5)
        
        # Home icon
        home = slide.shapes.add_textbox(prs.slide_width/2 - Inches(0.2), footer_y, Inches(0.4), Inches(0.4))
        home.text_frame.text = "🏠"
        home.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Left arrow
        left = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(0.4), Inches(0.4))
        left.text_frame.text = "◄"
        
        # Right arrow
        right = slide.shapes.add_textbox(prs.slide_width - Inches(0.9), footer_y, Inches(0.4), Inches(0.4))
        right.text_frame.text = "►"

    # SLIDE 1: TITLE PAGE
    slide1_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide1_layout)
    
    # Placeholder for Background Photo
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 20, 40) # Dark Navyish
    bg.line.fill.background()
    
    # Text
    title_box = slide1.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(3.0), Inches(5), Inches(2))
    tf = title_box.text_frame
    p1 = tf.add_paragraph()
    p1.text = "Express"
    p1.font.bold = True
    p1.font.size = Pt(60)
    p1.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Shopping"
    p2.font.bold = True
    p2.font.size = Pt(60)
    p2.font.color.rgb = WHITE
    
    p3 = tf.add_paragraph()
    p3.text = "E-Commerce"
    p3.font.bold = True
    p3.font.size = Pt(60)
    p3.font.color.rgb = WHITE
    
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.8), Inches(6), Inches(0.5))
    stf = subtitle_box.text_frame
    p = stf.add_paragraph()
    p.text = "Minor Project  |  React.js + Node.js + MongoDB"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE

    # SLIDE 2: INDEX
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Index"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = ORANGE_RED
    
    # List
    list_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(5))
    tf = list_box.text_frame
    items = ["Introduction", "Scope of Proposed System", "Objective of Proposed System", 
             "Task Dependency Diagram", "DFD / UML", "Step by Step Screen Layouts"]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = DEEP_NAVY
        p.space_after = Pt(12)
        # Bullet (simulated)
        p.text = "● " + item
        
    # Placeholder for image on right
    img_placeholder = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, prs.slide_width/2 + Inches(1), Inches(1.5), Inches(5), Inches(4))
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LAVENDER
    img_placeholder.line.color.rgb = DEEP_NAVY
    add_footer(slide2)

    # SLIDE 3: Introduction Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    # Fill background with placeholder for photo
    bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(100, 100, 150)
    
    title = slide3.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(1.5), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Introduction"
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    
    # Floating Lavender Cards on right
    card_w = Inches(3.5)
    card_h = Inches(2.0)
    cards = [
        {"title": "Easy Shopping", "text": "Express Shopping lets customers browse, search, and purchase products online from any device at any time of the day", "x": prs.slide_width - Inches(7.5), "y": Inches(1)},
        {"title": "Secure Payments", "text": "Safe and secure order placement with OTP verification, order history, and real-time order tracking for all users", "x": prs.slide_width - Inches(3.8), "y": Inches(1)},
        {"title": "Efficient Management", "text": "Admin can manage products, categories, orders, and users from a central dashboard without any physical interaction", "x": prs.slide_width - Inches(7.5), "y": Inches(3.5), "w": Inches(7.2)}
    ]
    
    for c in cards:
        w = c.get("w", card_w)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c["x"], c["y"], w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LAVENDER
        shape.line.color.rgb = WHITE
        
        tb = slide3.shapes.add_textbox(c["x"] + Inches(0.2), c["y"] + Inches(0.2), w - Inches(0.4), card_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.add_paragraph()
        p1.text = c["title"]
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = DEEP_NAVY
        
        p2 = tf.add_paragraph()
        p2.text = c["text"]
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        
    add_footer(slide3)

    # Continue with all 29 slides...
    # (I'll implement the rest in sections to keep it manageable but I MUST include all 29)
    
    # SLIDE 4: Introduction Detail
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    # Left illustration placeholder
    slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(5), Inches(4.5)).fill.solid()
    
    # Right large lavender card
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, prs.slide_width/2, 0, prs.slide_width/2, prs.slide_height)
    card.fill.solid()
    card.fill.fore_color.rgb = LAVENDER
    card.line.fill.background()
    
    tb = slide4.shapes.add_textbox(prs.slide_width/2 + Inches(0.5), Inches(1), prs.slide_width/2 - Inches(1.0), prs.slide_height - Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Express Shopping is a full-stack e-commerce web application that allows customers to browse products by category, add items to cart, and place orders online with ease.\n\nIt supports multiple user roles — Guest, Registered Customer, and Admin — with secure JWT authentication and role-based access for protected routes and admin features.\n\nBuilt with React.js for the frontend, Node.js and Express.js for the backend REST API, and MongoDB as the database."
    p.font.size = Pt(20)
    p.font.color.rgb = DEEP_NAVY
    add_footer(slide4)

    # SLIDE 5: Key Features
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Key Features"
    p.font.bold = True
    p.font.size = Pt(38)
    p.font.color.rgb = DEEP_NAVY
    
    cols = [
        {"title": "Product Shopping Cart", "text": "Add, remove, and update product quantities in cart. Persistent cart synced with user account via MongoDB.", "icon": "🛒"},
        {"title": "Secure Authentication", "text": "JWT-based login and registration. Role-based access control for Customer and Admin user types.", "icon": "🔐"},
        {"title": "Order Management", "text": "Place orders, view order history, and track order status from Pending to Shipped to Delivered.", "icon": "📦"}
    ]
    
    for i, col in enumerate(cols):
        x = Inches(0.5 + i * 4.2)
        # Icon
        ic = slide5.shapes.add_textbox(x, Inches(1.5), Inches(1), Inches(1))
        ic.text_frame.text = col["icon"]
        ic.text_frame.paragraphs[0].font.size = Pt(40)
        
        # Title
        tt = slide5.shapes.add_textbox(x, Inches(2.5), Inches(3.5), Inches(1))
        p = tt.text_frame.add_paragraph()
        p.text = col["title"]
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = DEEP_NAVY
        
        # Text
        tx = slide5.shapes.add_textbox(x, Inches(3.2), Inches(3.5), Inches(2))
        p = tx.text_frame.add_paragraph()
        p.text = col["text"]
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
    
    add_footer(slide5)

    # ... and so on ... I will group the rest into the file write
    
    # Slides 6-8: Scope
    for i in range(1, 4):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        p = title.text_frame.add_paragraph()
        p.text = "Scope Of Proposed System" + (f" Part {i}" if i > 1 else "")
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = DEEP_NAVY
        
        content = ""
        if i == 1:
            content = "✅ Online Product Browsing\n• Customers can browse products by category, brand, and price\n• Advanced search with keyword filtering and sort options\n• Product detail page with images, description, and reviews\n\n✅ Shopping Cart & Checkout\n• Add products to cart and update quantities anytime\n• Checkout with shipping address and payment method\n• Order summary and confirmation email after placing order"
        elif i == 2:
            content = "✅ User Account Management\n• Customer registration, login, and profile management\n• View order history, track order status, manage addresses\n• Secure password update and account deletion option\n\n✅ Admin Product Management\n• Add, edit, delete products with images and stock details\n• Manage product categories and featured product listings\n• View inventory levels and get low-stock alerts"
        else:
            content = "✅ Order & Payment Scope\n• View all customer orders with status update control\n• Filter orders by date, status, or customer name\n• Cash on Delivery and online payment method support\n\n✅ Scalability & Future Scope\n• Built with React.js and Node.js for high performance\n• Can be extended with mobile app, real payment, etc."
            
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        for line in content.split('\n'):
            p = tf.add_paragraph()
            p.text = line
            if '✅' in line:
                p.font.bold = True
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 150, 0)
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = GRAY
                p.level = 1 if '•' in line else 0
        
        # Right photo placeholder
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5)).fill.solid()
        add_footer(s)

    # Slides 9-11: Objectives
    objectives = [
        ["Secure User Authentication", "Implement JWT login and registration with bcrypt password hashing. Role-based access for Customer and Admin users."],
        ["Complete Product Catalog", "Enable customers to browse all products by category, search by keyword, sort by price, and filter by brand."],
        ["Smooth Cart & Checkout", "Allow users to add products to cart, update quantities, apply coupons, and place orders with shipping details."],
        ["RESTful API Backend", "Build Node.js + Express.js REST API with full CRUD operations, middleware, input validation, error handling."],
        ["MongoDB Database Design", "Design collections for Users, Products, Categories, Orders, Cart, Reviews using Mongoose ODM schemas."],
        ["Admin Management Panel", "Provide admin dashboard to manage products, orders, and customers with real-time stats and analytics."],
        ["Performance & Security", "Lazy loading React components, image optimization, Helmet.js headers, rate limiting, and CORS protection."],
        ["Cloud Deployment & CI/CD", "Deploy React frontend on Vercel, Node.js backend on Render, MongoDB Atlas for database, GitHub Actions CI/CD."]
    ]
    
    for slide_idx in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        p = title.text_frame.add_paragraph()
        p.text = "Objectives Of Proposed System"
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = DEEP_NAVY
        
        start = slide_idx * 3
        count = 3 if slide_idx < 2 else 2
        for j in range(count):
            idx = start + j
            y = Inches(1.5 + j * 1.8)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(6), Inches(1.6))
            card.fill.solid()
            card.fill.fore_color.rgb = LAVENDER
            card.line.color.rgb = DEEP_NAVY
            
            # Number badge
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.3), Inches(1), Inches(1))
            badge.fill.solid()
            badge.fill.fore_color.rgb = DEEP_NAVY
            bt = badge.text_frame
            p = bt.add_paragraph()
            p.text = str(idx + 1)
            p.font.bold = True
            p.font.size = Pt(30)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Text
            tx = s.shapes.add_textbox(Inches(1.8), y + Inches(0.2), Inches(4.5), Inches(1.2))
            tf = tx.text_frame
            p1 = tf.add_paragraph()
            p1.text = objectives[idx][0]
            p1.font.bold = True
            p1.font.size = Pt(18)
            p1.font.color.rgb = DEEP_NAVY
            
            p2 = tf.add_paragraph()
            p2.text = objectives[idx][1]
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY
            
        # Right photo
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Inches(5)).fill.solid()
        add_footer(s)

    # Slides 12-13: Diagrams
    for title_text in ["Task Dependency Diagram", "System Architecture"]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        p = title.text_frame.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = DEEP_NAVY
        # Diagram placeholder
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11), Inches(5)).fill.solid()
        add_footer(s)

    # Slides 14-16: DFD/UML
    for title_text in ["DFD / UML — Context Level DFD", "Level 1 DFD", "Use Case Diagram"]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        p = title.text_frame.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = DEEP_NAVY
        s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11), Inches(5)).fill.solid()
        add_footer(s)

    # Slide 17: UI Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "User Interface"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = DEEP_NAVY
    
    ui_items = ["Product Listing Page", "Simple & Clean Design", "Shopping Cart", "Customer Dashboard", 
                "Checkout Page", "Order Tracking", "Product Detail", "Admin Control Panel"]
    for i, item in enumerate(ui_items):
        col = i % 2
        row = i // 2
        x = Inches(1 + col * 6)
        y = Inches(1.5 + row * 1.2)
        tb = s.shapes.add_textbox(x, y, Inches(5), Inches(1))
        p = tb.text_frame.add_paragraph()
        p.text = "✨ " + item
        p.font.size = Pt(24)
        p.font.color.rgb = DEEP_NAVY
    add_footer(s)

    # Slides 18-28: Screen Layouts
    screens = [
        "Customer Login & Registration Page", "Customer Registration Page", "Home Page", 
        "Product Listing Page", "Product Detail Page", "Shopping Cart Page", 
        "Checkout Page", "Customer Dashboard & Profile", "Admin Login Page", 
        "Admin Dashboard", "Admin — Manage Products"
    ]
    for sc in screens:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(1))
        p = title.text_frame.add_paragraph()
        p.text = sc
        p.font.bold = True
        p.font.size = Pt(30)
        p.font.color.rgb = DEEP_NAVY
        # Large mockup placeholder
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.2), Inches(11.3), Inches(5.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = RGBColor(200, 200, 200)
        add_footer(s)

    # Slide 29: THANK YOU
    slide29 = prs.slides.add_slide(prs.slide_layouts[6])
    # Left BG Photo
    bg_left = slide29.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width/2, prs.slide_height)
    bg_left.fill.solid()
    bg_left.fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # Right Content
    title = slide29.shapes.add_textbox(prs.slide_width/2 + Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Thank You!"
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = DEEP_NAVY
    
    p2 = title.text_frame.add_paragraph()
    p2.text = "Thanks For Your Attention"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY
    
    # Navy Cards
    cards = [
        {"title": "Future Scope", "text": "Express Shopping can be extended with Razorpay/Stripe payment gateway, React Native mobile app, AI product recommendations, and real-time order tracking with GPS."},
        {"title": "Questions Welcome", "text": "Ready to discuss implementation details, architecture decisions, tech stack choices, and any customization."}
    ]
    for i, c in enumerate(cards):
        y = Inches(2.5 + i * 2.2)
        card = slide29.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, prs.slide_width/2 + Inches(0.5), y, Inches(5.5), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = DEEP_NAVY
        
        tb = slide29.shapes.add_textbox(prs.slide_width/2 + Inches(0.7), y + Inches(0.2), Inches(5.1), Inches(1.4))
        tf = tb.text_frame
        p1 = tf.add_paragraph()
        p1.text = c["title"]
        p1.font.bold = True
        p1.font.size = Pt(20)
        p1.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = c["text"]
        p2.font.size = Pt(14)
        p2.font.color.rgb = WHITE
        
    add_footer(slide29)

    prs.save('Express_Shopping_Presentation.pptx')
    print("Presentation created successfully.")

if __name__ == "__main__":
    create_presentation()
